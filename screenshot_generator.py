from playwright.sync_api import sync_playwright
from pptx import Presentation
from pptx.util import Inches
import os
import time

# 현재 디렉토리의 HTML 파일 경로
html_file = f"file:///{os.path.abspath('apm_price.html').replace(os.sep, '/')}"
output_image = "slide1.png"
viewport_width = 1920
viewport_height = 1080

print(f"HTML 파일 경로: {html_file}")

with sync_playwright() as p:
    browser = p.chromium.launch(headless=False)  # 브라우저 창 보이게
    page = browser.new_page(viewport={'width': viewport_width, 'height': viewport_height})
    
    print("페이지 로딩 중...")
    page.goto(html_file, wait_until='networkidle')
    
    # 폰트와 스타일이 완전히 로드될 때까지 대기
    time.sleep(3)
    
    print("스크린샷 캡처 중...")
    page.screenshot(path=output_image, full_page=True)  # 전체 페이지 캡처
    
    print(f"스크린샷 저장됨: {output_image}")
    browser.close()

# PPT 생성 및 이미지 삽입
print("PPT 생성 중...")
prs = Presentation()
blank_slide_layout = prs.slide_layouts[6]  # 빈 슬라이드

slide = prs.slides.add_slide(blank_slide_layout)

# 이미지 크기 정보 가져오기
from PIL import Image
img = Image.open(output_image)
img_width, img_height = img.size
img_ratio = img_width / img_height

# PPT 슬라이드 크기 (16:9 비율)
slide_width = prs.slide_width
slide_height = prs.slide_height
slide_ratio = slide_width / slide_height

print(f"이미지 크기: {img_width}x{img_height} (비율: {img_ratio:.2f})")
print(f"슬라이드 크기: {slide_width}x{slide_height} (비율: {slide_ratio:.2f})")

# 비율을 유지하면서 슬라이드에 맞는 크기 계산
if img_ratio > slide_ratio:
    # 이미지가 더 넓음 - 너비를 슬라이드에 맞춤
    new_width = slide_width
    new_height = int(slide_width / img_ratio)
    left = 0
    top = (slide_height - new_height) // 2
else:
    # 이미지가 더 높음 - 높이를 슬라이드에 맞춤
    new_height = slide_height
    new_width = int(slide_height * img_ratio)
    left = (slide_width - new_width) // 2
    top = 0

print(f"최적화된 크기: {new_width}x{new_height}, 위치: ({left}, {top})")

slide.shapes.add_picture(output_image, left, top, width=new_width, height=new_height)

ppt_filename = "apm_price_presentation.pptx"
prs.save(ppt_filename)
print(f"PPT 저장됨: {ppt_filename}") 