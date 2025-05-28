from pptx import Presentation
from pptx.util import Inches
import os

def html_to_ppt():
    # HTML 파일 읽기
    html_file = "apm_price.html"
    
    print(f"HTML 파일 읽는 중: {html_file}")
    
    with open(html_file, 'r', encoding='utf-8') as f:
        html_content = f.read()
    
    # PPT 생성
    prs = Presentation()
    
    # 슬라이드 크기를 A4 용지에 맞게 설정 (A4 비율: 210x297mm)
    prs.slide_width = Inches(11.69)   # A4 가로 (297mm)
    prs.slide_height = Inches(8.27)   # A4 세로 (210mm)
    
    # 빈 슬라이드 레이아웃 사용
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    
    # HTML을 웹 객체로 삽입하는 방법
    # 임시 HTML 파일의 절대 경로 생성
    abs_html_path = os.path.abspath(html_file)
    file_url = f"file:///{abs_html_path.replace(os.sep, '/')}"
    
    print(f"HTML URL: {file_url}")
    
    # 웹 브라우저 객체 추가 (OLE 객체)
    try:
        # PowerPoint에서 웹 페이지를 삽입하는 방법
        # 이것은 실제로는 iframe과 같은 역할을 함
        
        # 대신 HTML을 이미지로 변환해서 삽입하는 방식 사용
        from playwright.sync_api import sync_playwright
        
        print("HTML을 고품질 이미지로 변환 중...")
        
        with sync_playwright() as p:
            browser = p.chromium.launch(headless=True)
            # A4 비율에 맞는 뷰포트 설정 (1169x827 픽셀, A4 비율 유지)
            page = browser.new_page(viewport={'width': 1169, 'height': 827})
            
            # HTML 파일 로드
            page.goto(file_url, wait_until='networkidle')
            
            # 폰트와 스타일 완전 로딩 대기
            page.wait_for_timeout(3000)
            
            # 페이지의 실제 높이 확인
            page_height = page.evaluate("document.body.scrollHeight")
            print(f"페이지 실제 높이: {page_height}px")
            
            # 뷰포트를 실제 페이지 높이에 맞게 조정 (A4 비율 유지)
            if page_height > 827:
                # A4 비율을 유지하면서 높이 조정
                new_width = int(page_height * (1169/827))  # A4 비율 유지
                page.set_viewport_size({'width': new_width, 'height': page_height})
                page.wait_for_timeout(1000)  # 리사이즈 후 대기
            
            # 전체 페이지 스크린샷 생성
            temp_image = "temp_slide.png"
            page.screenshot(path=temp_image, full_page=True)
            browser.close()
        
        # 이미지를 슬라이드에 삽입 (비율 유지)
        from PIL import Image
        img = Image.open(temp_image)
        img_width, img_height = img.size
        img_ratio = img_width / img_height
        img.close()  # PIL 이미지 객체 닫기
        
        # PPT 슬라이드 크기
        slide_width = prs.slide_width
        slide_height = prs.slide_height
        slide_ratio = slide_width / slide_height
        
        print(f"이미지 크기: {img_width}x{img_height} (비율: {img_ratio:.2f})")
        print(f"슬라이드 크기: {slide_width}x{slide_height} (비율: {slide_ratio:.2f})")
        
        # 비율을 유지하면서 슬라이드를 꽉 채우도록 크기 계산 (crop 방식)
        if img_ratio > slide_ratio:
            # 이미지가 더 넓음 - 높이를 슬라이드에 맞추고 좌우는 잘림
            new_height = slide_height
            new_width = int(slide_height * img_ratio)
            left = (slide_width - new_width) // 2  # 중앙 정렬
            top = 0
        else:
            # 이미지가 더 높음 - 너비를 슬라이드에 맞추고 상하는 잘림
            new_width = slide_width
            new_height = int(slide_width / img_ratio)
            left = 0
            top = (slide_height - new_height) // 2  # 중앙 정렬
        
        print(f"최적화된 크기: {new_width}x{new_height}, 위치: ({left}, {top})")
        
        slide.shapes.add_picture(temp_image, left, top, width=new_width, height=new_height)
        
        # 임시 이미지 파일 삭제
        try:
            if os.path.exists(temp_image):
                os.remove(temp_image)
        except:
            print(f"임시 파일 {temp_image} 삭제 실패 (무시됨)")
            
        print("HTML 스타일이 완벽하게 보존된 PPT 생성 완료")
        
    except Exception as e:
        print(f"웹 객체 삽입 실패: {e}")
        print("대체 방법으로 텍스트 기반 슬라이드 생성...")
        
        # 기존 텍스트 기반 방법으로 폴백
        from pptx.util import Pt
        from pptx.enum.text import PP_ALIGN
        from pptx.dml.color import RGBColor
        
        # 제목 추가
        title_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(11.33), Inches(1))
        title_frame = title_box.text_frame
        title_p = title_frame.paragraphs[0]
        title_p.text = "모니터링 솔루션 공급 견적서"
        title_p.alignment = PP_ALIGN.CENTER
        title_p.font.size = Pt(36)
        title_p.font.bold = True
    
    # PPT 저장
    ppt_filename = "apm_price_direct.pptx"
    prs.save(ppt_filename)
    print(f"PPT 저장됨: {ppt_filename}")
    
    return ppt_filename

if __name__ == "__main__":
    html_to_ppt() 