from pptx import Presentation
from pptx.util import Inches
import os

def html_to_ppt_optimized():
    # A4 최적화 HTML 파일 사용
    html_file = "apm_price_a4.html"
    
    print(f"A4 최적화 HTML 파일 읽는 중: {html_file}")
    
    # PPT 생성
    prs = Presentation()
    
    # A4 용지 크기 설정 (정확한 A4 비율)
    prs.slide_width = Inches(11.69)   # A4 가로 (297mm)
    prs.slide_height = Inches(8.27)   # A4 세로 (210mm)
    
    # 빈 슬라이드 레이아웃 사용
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    
    # HTML 파일의 절대 경로 생성
    abs_html_path = os.path.abspath(html_file)
    file_url = f"file:///{abs_html_path.replace(os.sep, '/')}"
    
    print(f"HTML URL: {file_url}")
    
    try:
        from playwright.sync_api import sync_playwright
        
        print("고해상도 A4 최적화 이미지 생성 중...")
        
        with sync_playwright() as p:
            browser = p.chromium.launch(headless=True)
            
            # 고해상도 캡처를 위한 큰 뷰포트 (A4 비율 유지, 3배 확대)
            viewport_width = 1169 * 3  # 3507px
            viewport_height = 827 * 3   # 2481px
            
            page = browser.new_page(viewport={'width': viewport_width, 'height': viewport_height})
            
            # HTML 파일 로드
            page.goto(file_url, wait_until='networkidle')
            
            # 폰트와 스타일 완전 로딩 대기
            page.wait_for_timeout(3000)
            
            print(f"고해상도 뷰포트: {viewport_width}x{viewport_height}")
            
            # 고해상도 스크린샷 생성
            temp_image = "temp_slide_hd.png"
            page.screenshot(path=temp_image, full_page=False)
            browser.close()
        
        # 이미지를 슬라이드에 삽입 (전체 크기에 맞춤)
        slide.shapes.add_picture(temp_image, Inches(0), Inches(0), 
                                width=prs.slide_width, height=prs.slide_height)
        
        # 임시 이미지 파일 삭제
        try:
            if os.path.exists(temp_image):
                os.remove(temp_image)
        except:
            print(f"임시 파일 {temp_image} 삭제 실패 (무시됨)")
            
        print("A4 최적화 고해상도 PPT 생성 완료")
        
    except Exception as e:
        print(f"고해상도 캡처 실패: {e}")
        return None
    
    # PPT 저장
    ppt_filename = "apm_price_a4_fixed.pptx"
    prs.save(ppt_filename)
    print(f"A4 최적화 PPT 저장됨: {ppt_filename}")
    
    return ppt_filename

if __name__ == "__main__":
    html_to_ppt_optimized() 